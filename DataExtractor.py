"""
A concrete DataExtractor that fetches objects from a Google Cloud Storage
bucket, extracts plain‑text, and returns it to the GenAI assistant.

Supported file types:
    • PDF   → pdfminer.six
    • DOCX  → python-docx
    • PPTX  → python-pptx
    • TXT   → read as plain UTF‑8 text
    • Audio (wav, mp3, flac, m4a) → Cloud Speech‑to‑Text
    • Video (mp4, avi, mov, mkv) → Cloud Video Intelligence (speech transcription)

If an extraction fails, the method returns an empty string and prints a warning –
this keeps the rest of the pipeline running.
"""
cat > data_extractor.py << 'EOF'
import os
import tempfile
from typing import Optional

from google.cloud import storage, speech, videointelligence

# ---------- PDF ----------
from pdfminer.high_level import extract_text as pdf_extract_text

# ---------- DOCX ----------
import docx

# ---------- PPTX ----------
from pptx import Presentation


class DataExtractor:
    """
    Helper that knows how to pull a file from GCS and turn it into plain‑text.
    The constructor only stores the project/processor IDs – you can use them
    later if you want to call Document AI processors instead of the simple
    built‑in extractors.
    """

    def __init__(self, project_id: str, processor_id: str):
        self.project_id = project_id
        self.processor_id = processor_id

        # Clients are cheap to create; we keep a single instance per extractor.
        self.storage_client = storage.Client()
        self.speech_client  = speech.SpeechClient()
        self.video_client   = videointelligence.VideoIntelligenceServiceClient()

    # ------------------------------------------------------------------
    # Utility – download a GCS URI into a temporary local file and return its path.
    # The caller must delete the temporary file when done (we use `with` blocks).
    # ------------------------------------------------------------------
    def _download_to_temp(self, gcs_uri: str) -> str:
        """
        Returns the absolute path of a temporary file that contains the
        downloaded GCS object.  Caller should use `with tempfile.NamedTemporaryFile(...) as tmp:`
        to guarantee cleanup.
        """
        if not gcs_uri.startswith("gs://"):
            raise ValueError(f"Invalid GCS URI: {gcs_uri}")

        # Parse bucket & object name
        _, _, bucket_name, *object_parts = gcs_uri.split("/")
        object_name = "/".join(object_parts)

        bucket = self.storage_client.bucket(bucket_name)
        blob   = bucket.blob(object_name)

        # Create a NamedTemporaryFile that will be manually closed later
        tmp_file = tempfile.NamedTemporaryFile(delete=False)
        try:
            blob.download_to_filename(tmp_file.name)
        finally:
            tmp_file.close()                # keep the file on disk
        return tmp_file.name

    # ------------------------------------------------------------------
    # 1️⃣ Audio → Speech‑to‑Text
    # ------------------------------------------------------------------
    def extract_from_audio(self, gcs_uri: str) -> str:
        """Transcribe an audio file stored in GCS."""
        try:
            # Guess encoding from extension – fallback to UNSPECIFIED
            ext = os.path.splitext(gcs_uri)[1].lower()
            encoding_map = {
                ".wav": speech.RecognitionConfig.AudioEncoding.LINEAR16,
                ".flac": speech.RecognitionConfig.AudioEncoding.FLAC,
                ".mp3": speech.RecognitionConfig.AudioEncoding.MP3,
                ".m4a": speech.RecognitionConfig.AudioEncoding.MP4,
            }
            encoding = encoding_map.get(
                ext, speech.RecognitionConfig.AudioEncoding.ENCODING_UNSPECIFIED
            )

            audio = speech.RecognitionAudio(uri=gcs_uri)
            config = speech.RecognitionConfig(
                encoding=encoding,
                language_code="en-US",
                enable_automatic_punctuation=True,
            )
            response = self.speech_client.recognize(config=config, audio=audio)

            transcripts = [
                result.alternatives[0].transcript for result in response.results
            ]
            return " ".join(transcripts).strip()
        except Exception as e:
            print(f"[DataExtractor] Audio transcription failed for {gcs_uri}: {e}")
            return ""

    # ------------------------------------------------------------------
    # 2️⃣ Video → Speech transcription (Video Intelligence API)
    # ------------------------------------------------------------------
    def extract_from_video(self, gcs_uri: str) -> str:
        """Transcribe spoken words from a video stored in GCS."""
        try:
            features = [videointelligence.Feature.SPEECH_TRANSCRIPTION]
            operation = self.video_client.annotate_video(
                request={"features": features, "input_uri": gcs_uri}
            )
            result = operation.result(timeout=300)

            transcripts = []
            for annotation_result in result.annotation_results:
                for speech_transcription in annotation_result.speech_transcriptions:
                    best_alt = speech_transcription.alternatives[0]
                    transcripts.append(best_alt.transcript)

            return " ".join(transcripts).strip()
        except Exception as e:
            print(f"[DataExtractor] Video transcription failed for {gcs_uri}: {e}")
            return ""

    # ------------------------------------------------------------------
    # 3️⃣ PDF → plain text (pdfminer.six)
    # ------------------------------------------------------------------
    def extract_from_pdf(self, local_path: str) -> str:
        """Read a PDF file (already downloaded locally) and return its text."""
        try:
            return pdf_extract_text(local_path).strip()
        except Exception as e:
            print(f"[DataExtractor] PDF extraction failed for {local_path}: {e}")
            return ""

    # ------------------------------------------------------------------
    # 4️⃣ DOCX → plain text (python-docx)
    # ------------------------------------------------------------------
    def extract_from_docx(self, local_path: str) -> str:
        """Read a .docx file and return its text."""
        try:
            doc = docx.Document(local_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs).strip()
        except Exception as e:
            print(f"[DataExtractor] DOCX extraction failed for {local_path}: {e}")
            return ""

    # ------------------------------------------------------------------
    # 5️⃣ PPTX → plain text (python-pptx)
    # ------------------------------------------------------------------
    def extract_from_pptx(self, local_path: str) -> str:
        """Read a .pptx file and return its text."""
        try:
            prs = Presentation(local_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return "\n".join(texts).strip()
        except Exception as e:
            print(f"[DataExtractor] PPTX extraction failed for {local_path}: {e}")
            return ""

    # ------------------------------------------------------------------
    # 6️⃣ TXT → plain text (just read the file)
    # ------------------------------------------------------------------
    def extract_from_txt(self, local_path: str) -> str:
        """Read a plain‑text .txt file (UTF‑8) and return its content."""
        try:
            with open(local_path, "r", encoding="utf-8") as f:
                return f.read().strip()
        except Exception as e:
            print(f"[DataExtractor] TXT extraction failed for {local_path}: {e}")
            return ""

    # ------------------------------------------------------------------
    # Public API – these signatures are exactly what GenAIDocumentAssistant
    # calls.  No dispatcher, no recursion.
    # ------------------------------------------------------------------
    # (Audio & video already accept a GCS URI; PDF/DOCX/PPTX/TXT accept a local path)

    # keep the method names unchanged – the body simply calls the implementation above
    # (the “_” helpers have been merged into the methods themselves)